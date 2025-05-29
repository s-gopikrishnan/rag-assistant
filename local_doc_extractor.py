import os
import fitz  # PyMuPDF
import docx
from pptx import Presentation
import pdfplumber
from pathlib import Path
import json
from dataclasses import dataclass
from typing import List, Dict, Any, Optional
import hashlib
from datetime import datetime

@dataclass
class DocumentChunk:
    content: str
    metadata: Dict[str, Any]
    chunk_id: str
    doc_type: str
    source_file: str

class LocalDocumentExtractor:
    def __init__(self):
        self.supported_extensions = {'.pdf', '.docx', '.pptx'}
        
    def extract_document(self, file_path: str) -> List[DocumentChunk]:
        """Extract content from document and return structured chunks"""
        file_path = Path(file_path)
        
        if not file_path.exists():
            raise FileNotFoundError(f"File not found: {file_path}")
            
        extension = file_path.suffix.lower()
        
        if extension not in self.supported_extensions:
            raise ValueError(f"Unsupported file type: {extension}")
            
        # Route to appropriate extractor
        if extension == '.pdf':
            return self._extract_pdf(file_path)
        elif extension == '.docx':
            return self._extract_docx(file_path)
        elif extension == '.pptx':
            return self._extract_pptx(file_path)
    
    def _extract_pdf(self, file_path: Path) -> List[DocumentChunk]:
        """Extract content from PDF files"""
        chunks = []
        
        # Method 1: PyMuPDF for general extraction
        doc = fitz.open(file_path)
        
        # Extract document metadata
        doc_metadata = doc.metadata
        base_metadata = {
            'source_file': str(file_path),
            'doc_type': 'pdf',
            'total_pages': len(doc),
            'title': doc_metadata.get('title', ''),
            'author': doc_metadata.get('author', ''),
            'creation_date': doc_metadata.get('creationDate', ''),
            'extraction_timestamp': datetime.now().isoformat()
        }
        
        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            
            # Extract text
            text = page.get_text()
            
            if text.strip():  # Only process pages with content
                # Extract additional page metadata
                page_metadata = base_metadata.copy()
                page_metadata.update({
                    'page_number': page_num + 1,
                    'char_count': len(text),
                    'word_count': len(text.split())
                })
                
                # Create chunk ID
                chunk_id = self._generate_chunk_id(file_path, page_num, text[:100])
                
                chunk = DocumentChunk(
                    content=self._clean_text(text),
                    metadata=page_metadata,
                    chunk_id=chunk_id,
                    doc_type='pdf',
                    source_file=str(file_path)
                )
                chunks.append(chunk)
        
        doc.close()
        
        # Method 2: Use pdfplumber for tables if needed
        chunks.extend(self._extract_pdf_tables(file_path, base_metadata))
        
        return chunks
    
    def _extract_pdf_tables(self, file_path: Path, base_metadata: Dict) -> List[DocumentChunk]:
        """Extract tables from PDF using pdfplumber"""
        chunks = []
        
        with pdfplumber.open(file_path) as pdf:
            for page_num, page in enumerate(pdf.pages):
                tables = page.extract_tables()
                
                for table_idx, table in enumerate(tables):
                    if table and len(table) > 1:  # Has header and data
                        # Convert table to structured text
                        table_text = self._table_to_text(table)
                        
                        table_metadata = base_metadata.copy()
                        table_metadata.update({
                            'page_number': page_num + 1,
                            'content_type': 'table',
                            'table_index': table_idx,
                            'table_rows': len(table),
                            'table_cols': len(table[0]) if table else 0
                        })
                        
                        chunk_id = self._generate_chunk_id(file_path, f"table_{page_num}_{table_idx}", table_text[:100])
                        
                        chunk = DocumentChunk(
                            content=table_text,
                            metadata=table_metadata,
                            chunk_id=chunk_id,
                            doc_type='pdf',
                            source_file=str(file_path)
                        )
                        chunks.append(chunk)
        
        return chunks
    
    def _extract_docx(self, file_path: Path) -> List[DocumentChunk]:
        """Extract content from DOCX files"""
        chunks = []
        doc = docx.Document(file_path)
        
        # Extract document properties
        props = doc.core_properties
        base_metadata = {
            'source_file': str(file_path),
            'doc_type': 'docx',
            'title': props.title or '',
            'author': props.author or '',
            'created': props.created.isoformat() if props.created else '',
            'modified': props.modified.isoformat() if props.modified else '',
            'extraction_timestamp': datetime.now().isoformat()
        }
        
        # Extract paragraphs
        current_section = ""
        section_content = []
        
        for para_idx, paragraph in enumerate(doc.paragraphs):
            text = paragraph.text.strip()
            
            if not text:
                continue
                
            # Check if this is a heading (basic heuristic)
            is_heading = (
                paragraph.style.name.startswith('Heading') or
                len(text) < 100 and text.isupper() or
                (len(text.split()) < 10 and paragraph.runs and paragraph.runs[0].bold)
            )
            
            if is_heading and section_content:
                # Save previous section
                chunk = self._create_docx_chunk(
                    '\n'.join(section_content), 
                    current_section, 
                    base_metadata, 
                    file_path,
                    para_idx - len(section_content)
                )
                chunks.append(chunk)
                
                # Start new section
                current_section = text
                section_content = []
            else:
                section_content.append(text)
        
        # Add final section
        if section_content:
            chunk = self._create_docx_chunk(
                '\n'.join(section_content), 
                current_section, 
                base_metadata, 
                file_path,
                len(doc.paragraphs) - len(section_content)
            )
            chunks.append(chunk)
        
        # Extract tables
        for table_idx, table in enumerate(doc.tables):
            table_text = self._docx_table_to_text(table)
            
            table_metadata = base_metadata.copy()
            table_metadata.update({
                'content_type': 'table',
                'table_index': table_idx,
                'table_rows': len(table.rows),
                'table_cols': len(table.columns)
            })
            
            chunk_id = self._generate_chunk_id(file_path, f"table_{table_idx}", table_text[:100])
            
            chunk = DocumentChunk(
                content=table_text,
                metadata=table_metadata,
                chunk_id=chunk_id,
                doc_type='docx',
                source_file=str(file_path)
            )
            chunks.append(chunk)
        
        return chunks
    
    def _extract_pptx(self, file_path: Path) -> List[DocumentChunk]:
        """Extract content from PPTX files"""
        chunks = []
        prs = Presentation(file_path)
        
        # Extract presentation properties
        base_metadata = {
            'source_file': str(file_path),
            'doc_type': 'pptx',
            'total_slides': len(prs.slides),
            'extraction_timestamp': datetime.now().isoformat()
        }
        
        # Add title from first slide if available
        if prs.slides:
            first_slide = prs.slides[0]
            title_text = self._extract_slide_title(first_slide)
            if title_text:
                base_metadata['presentation_title'] = title_text
        
        for slide_idx, slide in enumerate(prs.slides):
            slide_content = []
            
            # Extract title
            title = self._extract_slide_title(slide)
            if title:
                slide_content.append(f"Slide Title: {title}")
            
            # Extract text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_content.append(shape.text.strip())
            
            if slide_content:
                content = '\n'.join(slide_content)
                
                slide_metadata = base_metadata.copy()
                slide_metadata.update({
                    'slide_number': slide_idx + 1,
                    'slide_title': title or f'Slide {slide_idx + 1}',
                    'char_count': len(content),
                    'word_count': len(content.split())
                })
                
                chunk_id = self._generate_chunk_id(file_path, slide_idx, content[:100])
                
                chunk = DocumentChunk(
                    content=self._clean_text(content),
                    metadata=slide_metadata,
                    chunk_id=chunk_id,
                    doc_type='pptx',
                    source_file=str(file_path)
                )
                chunks.append(chunk)
        
        return chunks
    
    def _extract_slide_title(self, slide) -> str:
        """Extract title from a slide"""
        if slide.shapes.title:
            return slide.shapes.title.text.strip()
        return ""
    
    def _create_docx_chunk(self, content: str, section_title: str, base_metadata: Dict, 
                          file_path: Path, para_start: int) -> DocumentChunk:
        """Create a chunk for DOCX content"""
        chunk_metadata = base_metadata.copy()
        chunk_metadata.update({
            'section_title': section_title,
            'paragraph_start': para_start,
            'char_count': len(content),
            'word_count': len(content.split())
        })
        
        chunk_id = self._generate_chunk_id(file_path, section_title, content[:100])
        
        return DocumentChunk(
            content=self._clean_text(content),
            metadata=chunk_metadata,
            chunk_id=chunk_id,
            doc_type='docx',
            source_file=str(file_path)
        )
    
    def _table_to_text(self, table: List[List[str]]) -> str:
        """Convert table data to structured text"""
        if not table:
            return ""
        
        # Assume first row is header
        headers = [cell or "" for cell in table[0]]
        rows = table[1:]
        
        text_parts = ["Table Content:"]
        text_parts.append("Headers: " + " | ".join(headers))
        
        for row_idx, row in enumerate(rows):
            row_data = [cell or "" for cell in row]
            text_parts.append(f"Row {row_idx + 1}: " + " | ".join(row_data))
        
        return "\n".join(text_parts)
    
    def _docx_table_to_text(self, table) -> str:
        """Convert DOCX table to text"""
        text_parts = ["Table Content:"]
        
        for row_idx, row in enumerate(table.rows):
            row_data = []
            for cell in row.cells:
                row_data.append(cell.text.strip())
            
            if row_idx == 0:
                text_parts.append("Headers: " + " | ".join(row_data))
            else:
                text_parts.append(f"Row {row_idx}: " + " | ".join(row_data))
        
        return "\n".join(text_parts)
    
    def _clean_text(self, text: str) -> str:
        """Clean extracted text"""
        # Remove excessive whitespace
        text = ' '.join(text.split())
        
        # Remove common artifacts
        text = text.replace('\x00', '')  # Null characters
        text = text.replace('\ufeff', '')  # BOM
        
        return text.strip()
    
    def _generate_chunk_id(self, file_path: Path, identifier: Any, content_preview: str) -> str:
        """Generate unique chunk ID"""
        source = f"{file_path.name}_{identifier}_{content_preview}"
        return hashlib.md5(source.encode()).hexdigest()[:16]

# Usage example
def process_document_folder(folder_path: str, output_path: str = None):
    """Process all documents in a folder"""
    extractor = LocalDocumentExtractor()
    all_chunks = []
    
    folder = Path(folder_path)
    for file_path in folder.rglob('*'):
        if file_path.suffix.lower() in extractor.supported_extensions:
            try:
                print(f"Processing: {file_path}")
                chunks = extractor.extract_document(file_path)
                all_chunks.extend(chunks)
                print(f"Extracted {len(chunks)} chunks from {file_path.name}")
            except Exception as e:
                print(f"Error processing {file_path}: {e}")
    
    # Optionally save to JSON for inspection
    if output_path:
        output_data = []
        for chunk in all_chunks:
            output_data.append({
                'chunk_id': chunk.chunk_id,
                'content': chunk.content,
                'metadata': chunk.metadata,
                'doc_type': chunk.doc_type,
                'source_file': chunk.source_file
            })
        
        with open(output_path, 'w', encoding='utf-8') as f:
            json.dump(output_data, f, indent=2, ensure_ascii=False)
    
    return all_chunks

# Example usage
if __name__ == "__main__":
    # Process single document
    extractor = LocalDocumentExtractor()
    chunks = extractor.extract_document("example.pdf")
    
    for chunk in chunks[:3]:  # Show first 3 chunks
        print(f"Chunk ID: {chunk.chunk_id}")
        print(f"Content: {chunk.content[:200]}...")
        print(f"Metadata: {chunk.metadata}")
        print("-" * 50)
    
    # Process entire folder
    # all_chunks = process_document_folder("./documents", "./extracted_data.json")